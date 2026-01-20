#!/usr/bin/env python3
"""
PowerPoint Generator Script
HTML/MarkdownからPPTXファイルを生成
"""

import argparse
import json
import os
import re
import sys
from pathlib import Path
from datetime import datetime
from typing import List, Dict, Any

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    from bs4 import BeautifulSoup
    import markdown
    from PIL import Image
except ImportError as e:
    print(f"Required package not installed: {e}")
    print("Run: pip install python-pptx beautifulsoup4 markdown pillow")
    sys.exit(1)

class PowerPointGenerator:
    def __init__(self, config_path: str = None):
        """PowerPoint生成器の初期化"""
        self.config = self.load_config(config_path)
        self.template_dir = Path(__file__).parent.parent / "templates"

    def load_config(self, config_path: str) -> Dict[str, Any]:
        """設定ファイルの読み込み"""
        default_config = {
            "defaultTemplate": "company-template.pptx",
            "branding": {
                "companyName": "株式会社CoCo",
                "logoPath": "./assets/logo.png",
                "primaryColor": "#007acc",
                "secondaryColor": "#6c757d"
            },
            "slideSettings": {
                "width": 13.33,
                "height": 7.5,
                "unit": "inches"
            }
        }

        if config_path and os.path.exists(config_path):
            with open(config_path, 'r', encoding='utf-8') as f:
                user_config = json.load(f)
                default_config.update(user_config)

        return default_config

    def parse_html_content(self, html_content: str) -> List[Dict[str, Any]]:
        """HTMLコンテンツをスライド形式にパース"""
        soup = BeautifulSoup(html_content, 'html.parser')
        slides = []

        # <section>タグでスライドを区切る場合
        sections = soup.find_all('section')
        if sections:
            for section in sections:
                slide_data = self.extract_slide_data(section)
                slides.append(slide_data)
        else:
            # コメントまたは見出しレベルでスライドを区切る
            slides = self.split_by_headers(soup)

        return slides

    def extract_slide_data(self, element) -> Dict[str, Any]:
        """HTML要素からスライドデータを抽出"""
        title = ""
        content = []
        images = []

        # タイトル取得
        title_elem = element.find(['h1', 'h2', 'h3'])
        if title_elem:
            title = title_elem.get_text().strip()
            title_elem.decompose()  # タイトル要素を削除
        elif element.get('data-slide-title'):
            title = element.get('data-slide-title')

        # 画像取得
        img_tags = element.find_all('img')
        for img in img_tags:
            if img.get('src'):
                images.append({
                    'src': img.get('src'),
                    'alt': img.get('alt', ''),
                    'width': img.get('width'),
                    'height': img.get('height')
                })
                img.decompose()  # 画像要素を削除

        # テキストコンテンツ取得
        paragraphs = element.find_all(['p', 'ul', 'ol', 'div', 'table'])
        for p in paragraphs:
            content.append({
                'type': p.name,
                'text': p.get_text().strip(),
                'html': str(p)
            })

        return {
            'title': title,
            'content': content,
            'images': images
        }

    def split_by_headers(self, soup) -> List[Dict[str, Any]]:
        """見出しレベルでHTMLを分割してスライド化"""
        slides = []
        current_slide = {'title': '', 'content': [], 'images': []}

        for element in soup.find_all(['h1', 'h2', 'h3', 'p', 'ul', 'ol', 'img', 'div']):
            if element.name in ['h1', 'h2', 'h3']:
                # 新しいスライドの開始
                if current_slide['title'] or current_slide['content']:
                    slides.append(current_slide)
                    current_slide = {'title': '', 'content': [], 'images': []}
                current_slide['title'] = element.get_text().strip()
            elif element.name == 'img':
                current_slide['images'].append({
                    'src': element.get('src'),
                    'alt': element.get('alt', ''),
                    'width': element.get('width'),
                    'height': element.get('height')
                })
            else:
                current_slide['content'].append({
                    'type': element.name,
                    'text': element.get_text().strip(),
                    'html': str(element)
                })

        # 最後のスライドを追加
        if current_slide['title'] or current_slide['content']:
            slides.append(current_slide)

        return slides

    def markdown_to_html(self, markdown_content: str) -> str:
        """MarkdownをHTMLに変換"""
        # スライド区切り（---）でセクション分割
        sections = markdown_content.split('---')
        html_sections = []

        for section in sections:
            section = section.strip()
            if section:
                html = markdown.markdown(section, extensions=['tables', 'fenced_code'])
                html_sections.append(f'<section>{html}</section>')

        return '\\n'.join(html_sections)

    def create_presentation(self, slides_data: List[Dict[str, Any]], template_path: str = None) -> Presentation:
        """スライドデータからPowerPointプレゼンテーションを作成"""
        if template_path and os.path.exists(template_path):
            prs = Presentation(template_path)
        else:
            # デフォルトテンプレートまたは新規作成
            default_template = self.template_dir / self.config['defaultTemplate']
            if default_template.exists():
                prs = Presentation(str(default_template))
            else:
                prs = Presentation()

        # 既存のスライドを削除（テンプレートの最初のスライド以外）
        while len(prs.slides) > 1:
            slide_to_remove = prs.slides[1]
            rId = prs.slides.index(slide_to_remove)
            prs.part.drop_rel(prs.slides._sldIdLst[rId].rId)
            del prs.slides._sldIdLst[rId]

        for slide_data in slides_data:
            self.add_slide_to_presentation(prs, slide_data)

        return prs

    def add_slide_to_presentation(self, prs: Presentation, slide_data: Dict[str, Any]):
        """プレゼンテーションにスライドを追加"""
        # レイアウトを選択（タイトルとコンテンツ）
        slide_layout = prs.slide_layouts[1] if len(prs.slide_layouts) > 1 else prs.slide_layouts[0]
        slide = prs.slides.add_slide(slide_layout)

        # タイトル設定
        if slide.shapes.title:
            slide.shapes.title.text = slide_data['title']
            self.format_title(slide.shapes.title)

        # コンテンツ追加
        content_placeholder = None
        for shape in slide.shapes:
            if shape.is_placeholder and shape.placeholder_format.idx == 1:
                content_placeholder = shape
                break

        if content_placeholder:
            self.add_content_to_placeholder(content_placeholder, slide_data['content'])

        # 画像追加
        self.add_images_to_slide(slide, slide_data['images'])

    def format_title(self, title_shape):
        """タイトルの書式設定"""
        if title_shape.has_text_frame:
            title_shape.text_frame.paragraphs[0].font.size = Pt(28)
            title_shape.text_frame.paragraphs[0].font.bold = True

            # プライマリカラーを適用
            primary_color = self.hex_to_rgb(self.config['branding']['primaryColor'])
            title_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(*primary_color)

    def add_content_to_placeholder(self, placeholder, content_list: List[Dict[str, Any]]):
        """プレースホルダーにコンテンツを追加"""
        if not placeholder.has_text_frame:
            return

        text_frame = placeholder.text_frame
        text_frame.clear()  # 既存のテキストをクリア

        for i, content in enumerate(content_list):
            if i > 0:
                p = text_frame.add_paragraph()
            else:
                p = text_frame.paragraphs[0]

            p.text = content['text']
            p.font.size = Pt(18)

            # リスト項目の場合はインデントを設定
            if content['type'] in ['ul', 'ol']:
                p.level = 1
                p.font.size = Pt(16)

    def add_images_to_slide(self, slide, images: List[Dict[str, Any]]):
        """スライドに画像を追加"""
        for i, img_data in enumerate(images):
            img_path = img_data['src']
            if os.path.exists(img_path):
                # 画像の配置位置を計算
                left = Inches(7) + (i * Inches(1))
                top = Inches(2)
                width = Inches(2)

                try:
                    slide.shapes.add_picture(img_path, left, top, width=width)
                except Exception as e:
                    print(f"Warning: Failed to add image {img_path}: {e}")

    def hex_to_rgb(self, hex_color: str) -> tuple:
        """HEX色をRGB色に変換"""
        hex_color = hex_color.lstrip('#')
        return tuple(int(hex_color[i:i+2], 16) for i in (0, 2, 4))

    def generate_from_file(self, input_path: str, output_path: str, template_path: str = None, input_format: str = 'auto'):
        """ファイルからPowerPointを生成"""
        print(f"[INFO] 入力ファイル読み込み: {input_path}")

        with open(input_path, 'r', encoding='utf-8') as f:
            content = f.read()

        # 入力形式を判定
        if input_format == 'auto':
            input_format = 'markdown' if input_path.endswith('.md') else 'html'

        # HTMLに変換
        if input_format == 'markdown':
            print("[INFO] Markdown → HTML変換")
            html_content = self.markdown_to_html(content)
        else:
            html_content = content

        # スライドデータ生成
        print("[INFO] スライドデータ解析")
        slides_data = self.parse_html_content(html_content)
        print(f"[INFO] {len(slides_data)}枚のスライドを生成")

        # PowerPoint作成
        print("[INFO] PowerPoint生成開始")
        prs = self.create_presentation(slides_data, template_path)

        # 保存
        prs.save(output_path)
        print(f"[INFO] PPTX出力完了: {output_path}")

def main():
    parser = argparse.ArgumentParser(description='HTML/MarkdownからPowerPointを生成')
    parser.add_argument('--input', '-i', required=True, help='入力ファイルパス')
    parser.add_argument('--output', '-o', required=True, help='出力PPTXファイルパス')
    parser.add_argument('--template', '-t', help='PowerPointテンプレートパス')
    parser.add_argument('--format', '-f', choices=['html', 'markdown', 'auto'], default='auto', help='入力形式')
    parser.add_argument('--config', '-c', help='設定ファイルパス')

    args = parser.parse_args()

    # 生成器作成
    generator = PowerPointGenerator(args.config)

    # PowerPoint生成
    try:
        generator.generate_from_file(args.input, args.output, args.template, args.format)
    except Exception as e:
        print(f"[ERROR] 生成エラー: {e}")
        sys.exit(1)

if __name__ == '__main__':
    main()